from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to set title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 136, 229) # Blue
        title.text_frame.paragraphs[0].font.bold = True

        # Content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default empty paragraph

        for item in content_text_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.space_after = Pt(14)
            p.level = 0

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    
    title.text = "高考模拟数据与志愿填报分析系统"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 136, 229)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "基于 Python Streamlit 的数据可视化解决方案\n\n汇报人：[你的名字]\n日期：2025年12月"

    # Slide 2: Project Background
    add_slide("项目背景与意义", [
        "📊 数据量大：高考数据繁杂，传统表格难以直观呈现。",
        "🤯 填报困难：考生难以快速定位自身位次与目标院校。",
        "💡 解决方案：构建一个可视化、交互式的分析看板。",
        "🎯 目标：实现从“查分”到“填报”的一站式辅助。"
    ])

    # Slide 3: 技术架构
    add_slide("技术栈与工具", [
        "🐍 Python 3.11：核心编程语言",
        "🎈 Streamlit：快速构建 Web 应用界面",
        "🐼 Pandas：强大的数据处理与清洗",
        "📈 Plotly Express：交互式数据可视化图表",
        "💻 VS Code + Copilot：高效开发环境"
    ])

    # Slide 4: 核心功能概览
    add_slide("系统核心功能", [
        "1️⃣ 成绩整体分析：全校/全班成绩分布、KPI指标。",
        "2️⃣ 个人成绩查询：详细成绩单、学科能力雷达图。",
        "3️⃣ 智能志愿推荐：基于“冲稳保”策略的院校推荐。",
        "4️⃣ 录取模拟演练：基于平行志愿算法的实时模拟。"
    ])

    # Slide 5: 功能展示 - 成绩分析
    add_slide("功能一：成绩整体分析", [
        "✅ 关键指标(KPI)：参考人数、平均分、最高/最低分。",
        "✅ 直方图：直观展示总成绩分布情况。",
        "✅ 箱线图：对比各学科成绩的离散程度。",
        "✅ 交互筛选：支持按班级筛选，实时更新图表。"
    ])

    # Slide 6: 功能展示 - 志愿推荐
    add_slide("功能二：智能志愿推荐", [
        "🔍 算法逻辑：基于“位次优先”与“分数线匹配”。",
        "📊 推荐策略：",
        "   - 冲：高于往年分数线 0-10 分",
        "   - 稳：高于往年分数线 10-30 分",
        "   - 保：高于往年分数线 30+ 分",
        "📋 结果展示：包含院校名称、最低投档分进度条。"
    ])

    # Slide 7: 功能展示 - 录取模拟
    add_slide("功能三：平行志愿录取模拟", [
        "⚙️ 核心算法：完全模拟真实的高考录取流程。",
        "🔄 流程：按位次排序 -> 检索6个志愿 -> 扣减名额。",
        "📂 结果输出：",
        "   - 实时显示录取/滑档人数。",
        "   - 支持下载 CSV 格式的详细录取名单。"
    ])

    # Slide 8: 总结与展望
    add_slide("总结与展望", [
        "✨ 成果：成功实现了一个功能完备的数据分析系统。",
        "🎨 体验：界面美观（自定义CSS），交互流畅。",
        "🚀 展望：",
        "   - 接入真实的高考历史数据。",
        "   - 增加更多维度的分析（如地区、专业倾向）。",
        "   - 引入 AI 大模型进行个性化咨询。"
    ])

    # Slide 9: Thank You
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "感谢观看"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 136, 229)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Q & A"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)

    prs.save('高考数据分析系统介绍.pptx')
    print("PPT generated successfully: 高考数据分析系统介绍.pptx")

if __name__ == "__main__":
    create_presentation()
